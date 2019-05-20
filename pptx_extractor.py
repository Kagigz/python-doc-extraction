import os
from pptx import Presentation

def extract_pptx(filename):

    document = Presentation(filename)
    ppt = {}

    # TITLE
    title = ""
    if(document.core_properties.title == ''):
        title = filename.split('.')[0]
    else:
        title = document.core_properties.title
    ppt['title'] = title

    # AUTHOR
    author = ""
    if(document.core_properties.author != ''):
        author = document.core_properties.author
    ppt['author'] = author

    # DATE
    date = ""
    if(document.core_properties.created != ''):
        date = document.core_properties.created
    ppt['date'] = date

    # MODIFIED DATE
    modDate = ""
    if(document.core_properties.modified != ''):
        modDate = document.core_properties.modified
    ppt["modDate"] = modDate


    # LANGUAGE
    lang = "fr"
    if(document.core_properties.language != ''):
        lang = document.core_properties.language
    ppt['lang'] = lang


    text = []
    for slide in document.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
    
    ppt['text'] = text

    return ppt


